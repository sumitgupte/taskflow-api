#!/usr/bin/env python3
"""Generate the facilitator slide deck from the bootcamp material.

The deck is facilitator-run: slides are deliberately sparse and the detail
lives in the speaker notes, which are lifted from docs/facilitator/README.md
and the four exercise handouts. Participants follow the handouts in the repo,
not these slides.

Usage:  python3 scripts/build-facilitator-deck.py [output.pptx]
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1A, 0x1F, 0x2B)
MUTED = RGBColor(0x5B, 0x64, 0x74)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WARN = RGBColor(0xC2, 0x41, 0x0C)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF4, 0xF6, 0xF9)
RULE = RGBColor(0xD8, 0xDE, 0xE6)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.85)
BODY_W = W - 2 * MARGIN

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _box(slide, left, top, width, height):
    from pptx.enum.shapes import MSO_SHAPE

    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, spacing=1.0):
    """runs: list of (text, size, bold, color, space_before_pt)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (text, size, bold, color, before) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(before)
        p.line_spacing = spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text.strip()


def title_slide(title, subtitle, footer, notes=""):
    s = prs.slides.add_slide(BLANK)
    _fill(_box(s, 0, 0, W, H), INK)
    _fill(_box(s, MARGIN, Inches(2.55), Inches(1.6), Inches(0.07)), ACCENT)
    _text(s, MARGIN, Inches(2.9), BODY_W, Inches(1.6),
          [(title, 46, True, PAPER, 0)])
    _text(s, MARGIN, Inches(4.25), BODY_W, Inches(1.0),
          [(subtitle, 21, False, RGBColor(0xB6, 0xC1, 0xD1), 0)])
    _text(s, MARGIN, Inches(6.45), BODY_W, Inches(0.5),
          [(footer, 13, False, RGBColor(0x7C, 0x8A, 0x9E), 0)])
    _notes(s, notes)
    return s


def section_slide(number, title, subtitle, notes=""):
    s = prs.slides.add_slide(BLANK)
    _fill(_box(s, 0, 0, W, H), WASH)
    _fill(_box(s, 0, 0, Inches(0.28), H), ACCENT)
    _text(s, MARGIN, Inches(2.5), BODY_W, Inches(0.7),
          [(number, 17, True, ACCENT, 0)])
    _text(s, MARGIN, Inches(3.1), BODY_W, Inches(1.3),
          [(title, 40, True, INK, 0)])
    _text(s, MARGIN, Inches(4.4), BODY_W, Inches(1.0),
          [(subtitle, 19, False, MUTED, 0)])
    _notes(s, notes)
    return s


def _header(s, title, eyebrow=None):
    top = Inches(0.62)
    if eyebrow:
        _text(s, MARGIN, top, BODY_W, Inches(0.35),
              [(eyebrow.upper(), 12, True, ACCENT, 0)])
        top = Inches(1.02)
    _text(s, MARGIN, top, BODY_W, Inches(0.9), [(title, 32, True, INK, 0)])
    _fill(_box(s, MARGIN, Inches(1.92), Inches(1.1), Inches(0.05)), ACCENT)


def bullets_slide(title, items, eyebrow=None, notes="", kicker=None):
    """items: list of str, or (str, sub) for a dimmed second line."""
    s = prs.slides.add_slide(BLANK)
    _fill(_box(s, 0, 0, W, H), PAPER)
    _header(s, title, eyebrow)
    runs = []
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, sub = item
            runs.append(("•  " + head, 20, True, INK, 0 if i == 0 else 17))
            runs.append(("     " + sub, 16, False, MUTED, 3))
        else:
            runs.append(("•  " + item, 20, False, INK, 0 if i == 0 else 15))
    _text(s, MARGIN, Inches(2.28), BODY_W, Inches(4.0), runs, spacing=1.08)
    if kicker:
        _fill(_box(s, MARGIN, Inches(6.35), BODY_W, Inches(0.62)), WASH)
        _text(s, MARGIN + Inches(0.25), Inches(6.5), BODY_W - Inches(0.5),
              Inches(0.4), [(kicker, 16, True, INK, 0)])
    _notes(s, notes)
    return s


def statement_slide(statement, attribution=None, notes="", dark=True):
    s = prs.slides.add_slide(BLANK)
    _fill(_box(s, 0, 0, W, H), INK if dark else WASH)
    fg = PAPER if dark else INK
    _text(s, MARGIN, Inches(2.6), BODY_W, Inches(2.4),
          [(statement, 34, True, fg, 0)], spacing=1.18)
    if attribution:
        _text(s, MARGIN, Inches(5.3), BODY_W, Inches(0.6),
              [(attribution, 17, False, RGBColor(0x9F, 0xAD, 0xC0) if dark else MUTED, 0)])
    _notes(s, notes)
    return s


def table_slide(title, headers, rows, eyebrow=None, notes="", widths=None):
    s = prs.slides.add_slide(BLANK)
    _fill(_box(s, 0, 0, W, H), PAPER)
    _header(s, title, eyebrow)
    n_rows, n_cols = len(rows) + 1, len(headers)
    top = Inches(2.28)
    height = Inches(min(4.4, 0.42 * n_rows))
    slide_table = s.shapes.add_table(n_rows, n_cols, MARGIN, top, BODY_W, height)
    tbl = slide_table.table
    tbl.first_row = True
    if widths:
        total = sum(widths)
        for i, w in enumerate(widths):
            tbl.columns[i].width = Emu(int(BODY_W * w / total))
    def style(cell, value, size, bold, color, bg):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left = cell.margin_right = Inches(0.12)
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        # An empty cell has no runs to style, so always add one explicitly.
        run = para.add_run()
        run.text = str(value)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

    for c, head in enumerate(headers):
        style(tbl.cell(0, c), head, 14, True, PAPER, INK)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            style(tbl.cell(r, c), val, 13, c == 0 and str(val) != "",
                  INK, PAPER if r % 2 else WASH)
    _notes(s, notes)
    return s


def compare_slide(title, left_title, left_items, right_title, right_items,
                  eyebrow=None, notes=""):
    s = prs.slides.add_slide(BLANK)
    _fill(_box(s, 0, 0, W, H), PAPER)
    _header(s, title, eyebrow)
    col_w = (BODY_W - Inches(0.6)) / 2
    for idx, (ct, items) in enumerate([(left_title, left_items),
                                       (right_title, right_items)]):
        left = MARGIN + idx * (col_w + Inches(0.6))
        _fill(_box(s, left, Inches(2.28), col_w, Inches(3.9)), WASH)
        _text(s, left + Inches(0.3), Inches(2.55), col_w - Inches(0.6),
              Inches(0.5), [(ct, 19, True, ACCENT if idx else WARN, 0)])
        runs = [("•  " + it, 15, False, INK, 0 if i == 0 else 11)
                for i, it in enumerate(items)]
        _text(s, left + Inches(0.3), Inches(3.15), col_w - Inches(0.6),
              Inches(3.0), runs, spacing=1.1)
    _notes(s, notes)
    return s


# ───────────────────────────────────────────────────────────── deck ──

title_slide(
    "Augmented AI Engineering",
    "Working with GitHub Copilot on a real codebase — a four-exercise workshop",
    "TaskFlow API  ·  Spring Boot 4 / Java 25  ·  2h25",
    """Housekeeping before you start:
- Confirm everyone has agent mode and a premium model. Exercises 2 and 3 compare two models; without a second model those steps collapse.
- java -version must say 25. Spring Boot 4.1 will not build on 17.
- Everyone should already have run ./mvnw test at home. First run pulls ~60MB from Maven Central — doing that on venue wifi at 09:01 with 20 people costs you the first exercise.
- Find out who is on IntelliJ vs VS Code. Both work; prompt files differ in Exercise 4.""")

statement_slide(
    "Everyone here can already get Copilot to write code.\n\nToday is about getting it to write the right code.",
    "The difference is a process, not a prompt.",
    """Set expectations in one minute. This is not a tips-and-tricks session and not a prompt-library handout.

The day is built as an experiment: the same feature gets built three times, varying one input at a time. That is the only way the comparison means anything.""")

bullets_slide(
    "How today works",
    [("One codebase", "TaskFlow API — small enough to read in fifteen minutes, real enough to break"),
     ("One ticket, built three times", "Exercises 2, 3 and 4 all build the same reminders feature"),
     ("One variable at a time", "Change the model, change the instructions, change the process"),
     ("You review everything", "Every exercise ends with reading generated code as if it were a PR")],
    eyebrow="Orientation",
    kicker="\"You're going to build this three times, and the third one is the only one you'd merge.\"",
    notes="""Say the kicker line out loud. If you do not, people experience Exercises 2–4 as repetition rather than as an experiment, and the whole arc is lost.

Each exercise branch is an independent snapshot off main. Nobody needs to have finished Exercise 1 to start Exercise 2 — a participant who falls behind or wrecks their tree switches branch and is caught up instantly.""")

table_slide(
    "Run of show",
    ["", "Block", "Mins", "Branch"],
    [["", "Intro — what Copilot is doing under the hood", "10", "—"],
     ["1", "Exploring a codebase with and without instructions", "20", "exercise-1"],
     ["", "Debrief 1", "5", ""],
     ["2", "Vibe code a feature without instructions", "20", "exercise-2"],
     ["", "Debrief 2", "5", ""],
     ["", "Break", "10", ""],
     ["3", "Implementing a feature ad-hoc with instructions", "30", "exercise-3"],
     ["", "Debrief 3", "5", ""],
     ["4", "Refine → plan → implement cycle", "35", "exercise-4"],
     ["", "Debrief 4 + close", "5", ""]],
    eyebrow="Agenda",
    widths=[0.06, 0.56, 0.12, 0.26],
    notes="""2h25 including the break.

The shape protects the back half: Exercises 3 and 4 keep 30 and 35 minutes, which is nearly the full original allowance. The compression falls on the intro, Exercise 1, Exercise 2, and the debriefs.

DEBRIEFS ARE FIVE MINUTES. That is two questions, not four. Each debrief slide marks the one question you must not skip — ask that, take two or three answers, move on. Overrunning here is the only way this shape breaks, because there is nowhere left to borrow from.

Debriefs are not optional padding — the arc only lands if you say the connecting sentence at each one.""")

bullets_slide(
    "The arc — say this out loud at the debriefs",
    [("Ex 1 — instructions change what Copilot knows, not just how it talks",
      "Behavioral vs structural instructions"),
     ("Ex 2 — no instructions, one-line ticket: output is a lottery",
      "A better model narrows the lottery. It is still a lottery."),
     ("Ex 3 — instructions and fitness functions raise the floor",
      "The ceiling does not move, because the ticket is still one line"),
     ("Ex 4 — the missing ingredient was never the prompt or the model",
      "It was that nobody had decided what to build")],
    eyebrow="Facilitator",
    notes="""Each exercise is engineered to fail in a way the next one fixes.

THE PIVOT AT THE END OF EXERCISE 3 IS THE WHOLE DAY. If people leave thinking the lesson is "write a good instructions file", the arc did not land.

Say the connecting sentence at every debrief or the sequence reads as four disconnected labs.""")

table_slide(
    "Branch model — each branch is a standalone snapshot",
    ["Branch", "Instructions", "OVERVIEW", "Guidelines", "Prompt files", "fitness.sh"],
    [["exercise-1", "✗", "✗", "✗", "✗", "✗"],
     ["exercise-2", "✗", "✗", "✗", "✗", "✗"],
     ["exercise-3", "✓", "✓", "✗ they write it", "✗", "✓"],
     ["exercise-4", "✓", "✓", "✓", "✓", "✓"],
     ["main", "✓", "✓", "✓", "✓ + facilitator guide", "✓"]],
    eyebrow="Setup",
    widths=[0.2, 0.16, 0.15, 0.2, 0.18, 0.13],
    notes="""Each branch carries only its own handout and its own tickets. This is deliberate: someone who reads the Exercise 4 handout while sitting in Exercise 2 already knows the punchline, and the day depends on Exercise 2 failing before Exercise 3 explains it.

fitness.sh is scoped the same way — Exercise 3 step 5 introduces it as the payoff of having just written guidelines, so shipping it earlier would pre-answer that discussion.

Reset within an exercise:  git checkout . && git clean -fd""")

bullets_slide(
    "What the short slot changed",
    [("Ex 2 — you run the repeat, on the projector",
      "Same prompt, same model, third time. Diff it live against your step-4 run."),
     ("Ex 4 step 5 — reading, not a second lab",
      "They open the other two tickets and discuss what Refine would find"),
     ("Debriefs — 5 minutes, two questions",
      "Each debrief slide marks the one you must not skip"),
     ("Ex 3 — untouched",
      "Keeps all 30 minutes and all four runs. It is the load-bearing comparison.")],
    eyebrow="Facilitator · 2h25 run",
    kicker="If you are behind entering Exercise 4, take it from step 6 — never from Refine.",
    notes="""One duty genuinely moves to you: the Exercise 2 repeat run.

Have your step-4 output still open and diff the two live — do not summarise it. They have to SEE two runs of an identical prompt disagreeing about what a reminder is. The handout tells them to expect this from you, so if you skip it the step just vanishes entirely.

Exercise 3 is deliberately untouched at 30 minutes. Its four runs are what produce the "model choice moved quality more than the instructions did" finding, and that finding is what makes the Exercise 4 pivot land. Do not borrow time from it.""")

# ── Intro block ──
section_slide("INTRO · 10 min", "What Copilot is actually doing",
              "Context in, tokens out — and everything you control is on the input side",
              """Keep this short and concrete. The goal is only to establish that the model sees a context window assembled from your prompt, your open files, whatever it retrieves, and any instructions files — and nothing else.

Everything for the rest of the day is a lever on that input.""")

compare_slide(
    "Two ways to change the answer",
    "Behavioral instructions", [
        "Change how Copilot answers",
        "Tone, length, format, what to avoid",
        "Add no new facts about your system",
        "Cost tokens on every single request",
        "\"Always respond concisely.\""],
    "Structural instructions", [
        "Change what Copilot knows",
        "Architecture, conventions, the files that matter",
        "Add facts it could not infer reliably",
        "The ones that change accuracy",
        "\"See PROJECT_OVERVIEW.md for the request flow.\""],
    eyebrow="Intro",
    notes="""People routinely conflate these two. Exercise 1 makes them feel the difference by asking the same question three times.

Do not pre-announce the answer to Exercise 1 here — set up the vocabulary only.""")

# ── Exercise 1 ──
section_slide("EXERCISE 1 · 20 min", "Exploring a codebase with and without instructions",
              "Branch: exercise-1 — no instructions file, no project overview",
              "Blank slate on purpose. Have everyone switch branch before you start talking.")

bullets_slide(
    "Ask the same question three times",
    [("1. Cold", "\"What does this code base do?\" — no instructions, no files attached. Record it."),
     ("2. After a behavioral rule", "Add one rule: \"Always respond concisely.\" New chat. Ask again."),
     ("3. After a structural one", "Generate PROJECT_OVERVIEW.md, review it, link it from the instructions. Ask again.")],
    eyebrow="Exercise 1 · steps",
    kicker="New chat each time — otherwise you are comparing against a primed conversation.",
    notes="""Expected shape of the three answers:

Step 2 (cold): a competent file-tree walk. Finds the controllers, correctly says "task management REST API". Typically misses that there is no real database, that the password is never checked, that NotificationService is an empty shell.

Step 4 (behavioral): noticeably shorter or longer. SAME FACTS, SAME GAPS. This is the point. If a table claims accuracy improved, ask them to point at the specific new fact.

Step 6 (structural): now cites the request flow and the limitations, because they are written down.""")

bullets_slide(
    "What Copilot invents in the generated overview",
    ["Claims login validates the password hash — it does not (see AuthController)",
     "Calls InMemoryDatabase \"a layer you could swap for JPA\" — no such abstraction exists",
     "Invents email or notification capability from the name NotificationService",
     "Misses that AuthInterceptor only guards /tasks/** — /health and /auth/login are open"],
    eyebrow="Exercise 1 · answer key",
    kicker="Plausible claims derived from names rather than code — the most useful thing to surface all day.",
    notes="""Push tables to find at least one of these. The User records do carry a bcrypt-looking hash field, which is what leads the model astray on the password claim.

If a table finishes early: have them add a WRONG fact to PROJECT_OVERVIEW.md on purpose and re-ask. Copilot repeats it confidently. Instructions are trusted input, not verified input.""")

bullets_slide(
    "Debrief 1",
    ["★  Which changed the answer more — the behavioral rule or the structural one?",
     "The overview was written by Copilot from code it already had. Why does feeding it back help?",
     "What did you have to correct in it? What does that say about an unreviewed overview?",
     "Behavioral rules cost tokens on every request. Which would you actually keep?"],
    eyebrow="Exercise 1 · debrief",
    kicker="Takeaway: structural instructions change the answers — and are only as good as your review.",
    notes="""FIVE MINUTES. Ask the starred question, take two or three answers, then pick ONE more. The others are in the handout for them to think about later.

Question 2 is the interesting one. The answer: retrieval is unreliable and lossy, and an instructions file is guaranteed context rather than hoped-for context. It is a cache of a review you already did.

Land the takeaway, then move on — do not let this run long, Exercise 2 needs its full 30.""")

# ── Exercise 2 ──
section_slide("EXERCISE 2 · 20 min", "Vibe code a feature without instructions",
              "Branch: exercise-2 — same repo, no instructions, one-line ticket",
              "Say the 'three times' line here if you have not already.")

statement_slide(
    "\"Users keep missing deadlines. Can we add reminders so they get\nnotified before a task is due? Should work for all tasks.\nLet's ship this soon.\"",
    "That is the entire ticket. This is realistic.",
    """Read it out loud. Then say: everything that goes wrong in the next thirty minutes traces back to something this ticket does not say.

Do not list the six ambiguities yet. Let them discover the failures first — the list is the debrief, not the setup.""")

bullets_slide(
    "Three runs, same feature — plus one you'll watch",
    [("Run 1 — GPT-5 mini", "The lazy path. Don't help it, don't clarify, don't correct mid-flight."),
     ("Run 2 — a premium model", "Same prompt, word for word. Note the differences specifically."),
     ("On the projector — the same prompt again", "Your facilitator runs it. Is the result the same?"),
     ("Run 3 — \"...using TDD\"", "Three extra words. Compare to all of the above.")],
    eyebrow="Exercise 2 · steps",
    kicker="Commit after every run:  git add -A && git commit -m \"ex2: gpt-5-mini, no instructions\"",
    notes="""PUSH PEOPLE TO COMMIT. Exercise 4 step 4 has them diff all three exercises against each other, and that payoff does not work if nobody committed.

Reset between runs:  git switch exercise-2 && git checkout . && git clean -fd

Step 5 (same prompt, same model, different result) is easy to rush past. Make people actually diff it. "It worked when I tried it" dies here.""")

bullets_slide(
    "What the runs typically produce",
    ["A reminderSent field bolted onto Task, with nothing that ever fires",
     "@Scheduled invented out of nowhere — sometimes a fabricated Quartz dependency",
     "An email sender hallucinated behind NotificationService, with SMTP config for no server",
     "NPE on the seeded null dueDate values — the cheapest proof it never read InMemoryDatabase",
     "Reminders fired for t3, which is already done",
     "GET /reminders unscoped by user, so Alice reads Bob's"],
    eyebrow="Exercise 2 · answer key",
    kicker="Steer at least one table to the unscoped endpoint — it maps straight onto a real incident.",
    notes="""Roughly in order of frequency. The class name NotificationService alone is enough to trigger the invented email transport.

CALIBRATION: this list was written against GPT-4.1 vs Sonnet 4.5. With GPT-5 mini at the weak end, expect fewer outright NPEs and more plausible-looking code that still answers the six ambiguities silently. That is a better demonstration, not a worse one — "it compiles and looks reasonable" is the failure mode people actually meet at work. Steer the review toward WHAT WAS DECIDED rather than what broke.

With Opus 5 at the strong end, it may ask a clarifying question or refuse to invent a scheduler. Name it when it happens: the model did the Refine stage on its own, partially, because nobody else had. Perfect setup for Exercise 4.

The unscoped read endpoint is the highest-value moment in the exercise. If no table finds it, ask directly: "can Alice see Bob's reminders?"

Watch for a model inventing a scheduler and never mentioning that it did.""")

bullets_slide(
    "The six questions every implementation answered",
    ["What is a reminder here — in-app, email, push?",
     "How long before the due date does it fire? Fixed, or per task?",
     "What happens to tasks with no due date?",
     "What happens to tasks already marked done?",
     "Does it fire once, or keep reminding until the task is done?",
     "Who can see or manage another user's reminders?"],
    eyebrow="Exercise 2 · debrief",
    kicker="The model didn't skip these decisions. It made them silently, and differently every run.",
    notes="""This is the scoring rubric for the rest of the day. Put it on the wall if you can.

The point to land: every implementation answers all six. The only question is whether a human decided the answers or a model guessed them.

On step 6 (latent knowledge): three words — "using TDD" — restructured the output. Nobody explained TDD to the model. The leverage was in NAMING A PRACTICE THE MODEL ALREADY KNOWS, not in describing it. That is the cheapest prompt improvement available.""")

bullets_slide(
    "Debrief 2",
    ["Across your attempts, how many produced a feature you would merge?",
     "★  Did any two runs agree on what \"reminder\" meant? If not, whose fault is that?",
     "Which was the bigger lever — changing the model, or changing the prompt?",
     "You reviewed generated code four times. How does that time compare to what generation saved?"],
    eyebrow="Exercise 2 · debrief",
    kicker="Takeaway: a better model narrows the variance. It does not remove it.",
    notes="""FIVE MINUTES. The starred question is the one that must land — everything else here is optional.

Question 2 is the trap worth springing: the instinct is to blame the model. The answer is that nobody told it, so it is the ticket's fault.

Close by naming the gap neither a better model nor a better prompt closed — the ticket never said what it wanted. That is what Exercises 3 and 4 attack.""")

# ── Exercise 3 ──
section_slide("EXERCISE 3 · 30 min", "Implementing a feature ad-hoc with instructions",
              "Branch: exercise-3 — same ticket, same models, instructions now in place",
              """Hold the feature constant, vary one input. If instructions matter, this is where you see it — and if they matter less than expected, you see that too.

Be honest about the result. It sets up Exercise 4.""")

bullets_slide(
    "Same ticket. What changed is what the repo tells Copilot first.",
    [(".github/copilot-instructions.md", "A behavioral rule plus a pointer to the overview"),
     ("PROJECT_OVERVIEW.md", "The reviewed architecture description from Exercise 1"),
     ("scripts/fitness.sh", "Tests + format check + compile — first appears on this branch"),
     ("docs/coding-guidelines.md", "Not here yet. They write it in step 5.")],
    eyebrow="Exercise 3 · setup",
    notes="""Have them read .github/copilot-instructions.md before they start, so they know what Copilot is being told that it was not told in Exercise 2.

Same four-run structure as the full-length version: weak model, then premium model, then write the guidelines, then a final run with auto-approve on. This exercise keeps its whole 30 minutes — it is the one place in the short slot where nothing was cut.

Shape: two runs with reviews (18 min), write the guidelines (6), final run with auto-approve on (6).""")

statement_slide(
    "Model choice moves quality more than the instructions file does\n— for a ticket this vague.",
    "Let that land honestly rather than steering to \"instructions are great\".",
    """Expected ranking of the four runs: Ex2 weak → Ex2 strong → Ex3 weak → Ex3 strong.

What instructions DO reliably improve: adherence to existing patterns (constructor injection, ErrorResponse), and a sharp drop in invented schedulers and email transports — because PROJECT_OVERVIEW.md says in as many words that neither exists. That is the clearest measurable win. Point at it.

What they do not fix: the feature still is not specified. The six ambiguities are still being answered silently, just in better-formatted code.

Have tables run  git diff exercise-2 -- src/  rather than eyeballing it.""")

bullets_slide(
    "Fitness functions — a guideline nobody can check is a wish",
    [("./scripts/fitness.sh", "Runs the test suite, the format check, and a compile"),
     ("The discussion is what it doesn't gate", "Read your own guidelines against it and find the gaps"),
     ("Scoping by userId", "→ an ArchUnit rule, or a test that logs in as Bob and asserts nothing of Alice's"),
     ("Changelog updated", "→ a CI check that the diff touches CHANGELOG.md"),
     ("No trivial comments", "→ not mechanically checkable. That is the honest answer.")],
    eyebrow="Exercise 3 · step 5",
    kicker="The rules with no automated check behind them are the ones that keep getting violated.",
    notes=""""Then it needs a human in review" is the right conclusion for the third one — do not let a table invent a fake linter for it.

Guidelines earn their place by being the write-up of a mistake you have already seen. The "don't invent infrastructure" rule exists because they watched it happen twice this morning.

On auto-approval (step 6): some participants will be nervous. That is the correct instinct and worth naming — auto-approve is defensible only in proportion to the strength of your automated gates. A repo with no fitness.sh has no business running an unsupervised agent.""")

bullets_slide(
    "Debrief 3",
    ["★  Rank all the runs across Exercises 2 and 3. What does the ordering say?",
     "Which guideline did the model violate anyway? What would catch it automatically?",
     "Auto-approval made it faster. Did it make it better?",
     "Your guidelines load on every request. What would you cut from a 400-line file first?"],
    eyebrow="Exercise 3 · debrief",
    kicker="Takeaway: instructions enforce your review comments before the code is written.",
    notes="""FIVE MINUTES. Ask the starred question — the ranking is what sets up the next slide, and the next slide is the whole day.

Expected ordering: Ex2 weak → Ex2 strong → Ex3 weak → Ex3 strong. Let the honest finding land: for a ticket this vague, model choice moved quality more than the instructions file did.

If they are slow to answer question 3, name the nervousness directly — auto-approve is defensible only in proportion to the strength of your automated gates. A repo with no fitness.sh has no business running an unsupervised agent.

Go straight into the pivot slide from here. Do not break the momentum.""")

statement_slide(
    "Six runs in, nobody has yet decided what a reminder is.",
    "Which is Exercise 4. — This pivot is the whole day.",
    """THIS IS THE MOST IMPORTANT SLIDE IN THE DECK.

If people leave thinking the lesson was "write a good instructions file", the arc did not land. Say explicitly: instructions turned your review comments into something enforced before the code is written rather than after. Fitness functions turned them into something enforced without you in the room.

Neither can tell Copilot what the feature is supposed to DO.

Take the break here if you have not already.""")

# ── Exercise 4 ──
section_slide("EXERCISE 4 · 35 min", "Refine → plan → implement",
              "Branch: exercise-4 — third time building the same feature",
              "Stop trying to fix it with a better prompt. Fix it with a better process.")

table_slide(
    "Three prompts, in order, each reviewed before the next runs",
    ["Stage", "Question it answers", "Output"],
    [["Refine", "What are we actually building?", "A spec — docs/tickets/reminders-spec.md"],
     ["Plan", "How, in what order, verified how?", "Numbered steps — reminders-plan.md"],
     ["Implement", "Do it.", "Code + tests, fitness.sh green"]],
    eyebrow="Exercise 4",
    widths=[0.18, 0.42, 0.4],
    notes="""Expected time split: Refine 13, Plan 9, Implement 13 — and they will not finish the feature. SAY SO UP FRONT, TWICE.

The failure mode is a table rushing Refine to "get to the real work", which is exactly the habit this exercise exists to break.

This branch ships all three as prompt files in .github/prompts/, invoked as /refine, /plan, /implement. On IntelliJ slash commands may not resolve — the handout gives every prompt in full under "…or by hand".""")

bullets_slide(
    "What good looks like at each stage",
    [("A good spec settles", "In-app only. Fixed lead time via properties. No dueDate → skipped. Done → skipped. One-shot. Owner-only."),
     ("A good plan", "First step is a test. No step you cannot verify. \"3. Implement the reminder service\" is the tell."),
     ("A good implement run", "States which step it is on. Runs fitness.sh rather than claiming it. Stops at the plan's boundary.")],
    eyebrow="Exercise 4 · answer key",
    notes="""Common implement-stage failures, all worth calling out:
- Claims ./scripts/fitness.sh is green without running it. Have them check.
- Drifts off the plan around step 3 and starts designing something else.
- Implements steps 1–5 in one turn, so nothing was verified independently.

Have pairs swap plans and review each other's. A step you cannot tell "done" for is not a step.""")

bullets_slide(
    "The comparison that makes the point",
    ["git diff exercise-2 -- src/     — vibe coded, no instructions",
     "git diff exercise-3 -- src/     — instructions and guidelines, still a one-line prompt",
     "Score all three on the six questions from this morning",
     "Exercise 2's code probably answers all six too"],
    eyebrow="Exercise 4 · step 4",
    kicker="The difference isn't how many got answered. It's that yours are written down and reviewed.",
    notes="""This is the payoff of the entire day, and it only works if people committed after each exercise. If a table did not commit, pair them with one that did.

Step 5 — the second ticket — is the other half of the lesson. task-listing-performance.md is fully scoped, so Refine finds almost nothing. That shows the cycle's value is proportional to the ticket's vagueness, not a ritual to perform on everything.

If time is short: half the room takes rate-limiting, half takes performance, then compare Refine outputs.""")

bullets_slide(
    "When is this cycle NOT worth it?",
    [("rate-limiting.md", "Already half-refined — Refine has less to find"),
     ("task-listing-performance.md", "Fully scoped, no ambiguity. You'd go straight to Plan."),
     ("The reminders ticket", "One line. Refine was the only thing that made the rest work.")],
    eyebrow="Exercise 4 · step 5",
    kicker="The value of the cycle scales with the vagueness of the ticket. It is not a ritual.",
    notes="""In this slot it is a READING step, not a lab — there is no time to run the cycle twice. They open both tickets and answer "what would Refine even find here?"

LAND THIS EXPLICITLY. It is the single most important thing for adoption back at their desks. A room that leaves believing every change needs a three-stage cycle will abandon the whole idea within a fortnight.

If a table finished the implement stage early, have them actually run the cycle on rate-limiting.md instead of just reading it.""")

bullets_slide(
    "Go further — what else belongs in the cycle?",
    [("/review", "Review this diff against the plan and the guidelines, as a PR"),
     ("/critique", "Argue against this plan. What's the strongest case it's wrong?"),
     ("/decompose", "Split step 4 into steps verifiable in under ten minutes"),
     ("/regress", "What existing behavior could this break? Write a test for each")],
    eyebrow="Exercise 4 · step 6",
    kicker="Plan adherence: one step per turn · check steps off in the plan file · make drift fail loudly",
    notes="""This is a discussion, not a lab.

The plan-adherence answers you are fishing for: number the steps and have it state which one it is on; have it check off completed steps in the plan file so the plan becomes state; implement one step per chat turn; reference the plan by file rather than pasting it; make the fitness functions strict enough that drifting fails rather than merely being wrong.""")

bullets_slide(
    "Debrief 4",
    ["★  Same ticket, three processes. How much is the process, and how much is just spending longer?",
     "Which of the three stages was hardest not to skip under time pressure?",
     "Refine took 15 minutes and produced no code. Sell that to a sceptical colleague in two sentences.",
     "Where's the cutoff? What size of change is not worth this cycle?"],
    eyebrow="Exercise 4 · debrief",
    notes="""FIVE MINUTES, and it is the last thing before the close — do not let it sprawl.

Question 1 is the fair challenge and you should not dodge it. Part of the answer is honestly "yes, you spent longer" — but the spec and plan are reusable artefacts and the review was cheaper because the diff was smaller and expected.

Question 4 matters for adoption: nobody should run refine→plan→implement on a typo fix. The cycle scales with the vagueness of the ticket.""")

statement_slide(
    "Refine decides what's true.\nPlan decides what order.\nImplement is the part everyone thinks is the work.",
    "The model is the same in all three stages. You're the thing that changed.",
    """Close on this. Then hand out the cheat sheet.

docs/copilot-prompt-cheatsheet.md is scoped to exercise-4 and has prompt templates for each stage. It is the thing people keep after the bootcamp.""")

bullets_slide(
    "What to take back",
    [("A PROJECT_OVERVIEW.md your team reviewed", "Structural context beats behavioral rules"),
     ("Coding guidelines that are write-ups of real mistakes", "Not a style wishlist"),
     ("A fitness.sh that actually gates something", "Auto-approval is only as safe as your gates"),
     ("Refine → plan → implement, on vague tickets only", "The cycle scales with the vagueness")],
    eyebrow="Close",
    kicker="docs/copilot-prompt-cheatsheet.md — keep this one.",
    notes="""Ask each table for one thing they will change on their own repo on Monday. Write them on the whiteboard. Concrete commitments beat a satisfaction score.""")

out = Path(sys.argv[1] if len(sys.argv) > 1
           else "docs/facilitator/bootcamp-deck.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Wrote {out} — {len(prs.slides._sldIdLst)} slides")
